"""Document processing service for extracting text from various file formats."""

import io
from pathlib import Path

from docx import Document as DocxDocument
from pptx import Presentation
from PyPDF2 import PdfReader

from ..core.logger import logging

logger = logging.getLogger(__name__)


class DocumentProcessor:
    """Service for processing and extracting text from documents."""

    @staticmethod
    async def extract_text_from_pdf(file_content: bytes) -> str:
        """Extract text from PDF file.

        Args:
            file_content: PDF file content as bytes

        Returns:
            Extracted text content

        Raises:
            Exception: If PDF processing fails
        """
        try:
            pdf_file = io.BytesIO(file_content)
            pdf_reader = PdfReader(pdf_file)

            text_content = []
            for page in pdf_reader.pages:
                text = page.extract_text()
                if text.strip():
                    text_content.append(text)

            extracted_text = "\n\n".join(text_content)
            logger.info(f"Successfully extracted {len(extracted_text)} characters from PDF")
            return extracted_text

        except Exception as e:
            logger.error(f"Error extracting text from PDF: {str(e)}")
            raise Exception(f"Failed to extract text from PDF: {str(e)}")

    @staticmethod
    async def extract_text_from_pptx(file_content: bytes) -> str:
        """Extract text from PowerPoint file.

        Args:
            file_content: PPTX file content as bytes

        Returns:
            Extracted text content

        Raises:
            Exception: If PPTX processing fails
        """
        try:
            pptx_file = io.BytesIO(file_content)
            presentation = Presentation(pptx_file)

            text_content = []
            for slide_num, slide in enumerate(presentation.slides, 1):
                slide_text = []

                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())

                if slide_text:
                    text_content.append(f"Slide {slide_num}:\n" + "\n".join(slide_text))

            extracted_text = "\n\n".join(text_content)
            logger.info(f"Successfully extracted {len(extracted_text)} characters from PPTX")
            return extracted_text

        except Exception as e:
            logger.error(f"Error extracting text from PPTX: {str(e)}")
            raise Exception(f"Failed to extract text from PPTX: {str(e)}")

    @staticmethod
    async def extract_text_from_docx(file_content: bytes) -> str:
        """Extract text from Word document.

        Args:
            file_content: DOCX file content as bytes

        Returns:
            Extracted text content

        Raises:
            Exception: If DOCX processing fails
        """
        try:
            docx_file = io.BytesIO(file_content)
            document = DocxDocument(docx_file)

            text_content = []
            for paragraph in document.paragraphs:
                if paragraph.text.strip():
                    text_content.append(paragraph.text.strip())

            extracted_text = "\n\n".join(text_content)
            logger.info(f"Successfully extracted {len(extracted_text)} characters from DOCX")
            return extracted_text

        except Exception as e:
            logger.error(f"Error extracting text from DOCX: {str(e)}")
            raise Exception(f"Failed to extract text from DOCX: {str(e)}")

    @staticmethod
    async def process_document(file_content: bytes, file_type: str) -> str:
        """Process document and extract text based on file type.

        Args:
            file_content: File content as bytes
            file_type: Type of file (pdf, pptx, docx)

        Returns:
            Extracted text content

        Raises:
            ValueError: If file type is not supported
            Exception: If document processing fails
        """
        file_type = file_type.lower()

        if file_type == "pdf":
            return await DocumentProcessor.extract_text_from_pdf(file_content)
        elif file_type == "pptx":
            return await DocumentProcessor.extract_text_from_pptx(file_content)
        elif file_type == "docx":
            return await DocumentProcessor.extract_text_from_docx(file_content)
        else:
            raise ValueError(f"Unsupported file type: {file_type}")

    @staticmethod
    def validate_file_type(filename: str) -> str:
        """Validate and extract file type from filename.

        Args:
            filename: Name of the file

        Returns:
            File extension/type

        Raises:
            ValueError: If file type is not supported
        """
        file_path = Path(filename)
        file_extension = file_path.suffix.lower().lstrip('.')

        supported_types = ["pdf", "pptx", "docx"]
        if file_extension not in supported_types:
            raise ValueError(f"Unsupported file type: {file_extension}. Supported types: {supported_types}")

        return file_extension
